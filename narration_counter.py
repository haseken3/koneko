"""ko-NeKo — PPTXナレーション文字数カウント＋尺推定ロジック"""

import re
from pptx import Presentation

# ポーズ時間（VoiceSpace デフォルト設定）
PAUSE_COMMA = 0.3   # 「、」1個あたり（秒）
PAUSE_PERIOD = 1.0  # 「。」1個あたり（秒）


def _count_reading_chars(text: str) -> int:
    """読み上げ対象の文字数をカウント（空白・改行を除外）。"""
    return len(re.sub(r'\s', '', text))


def _count_pauses(text: str) -> float:
    """句読点によるポーズ時間を算出（秒）。"""
    commas = text.count('、')
    periods = text.count('。')
    return commas * PAUSE_COMMA + periods * PAUSE_PERIOD


def extract_slide_notes(pptx_file):
    """全スライドのノート欄テキスト・文字数を抽出する。

    Args:
        pptx_file: PPTXファイルパスまたはファイルライクオブジェクト

    Returns:
        list[dict]: 各スライドの情報
            slide_num, title, notes, char_count, pause_seconds
    """
    prs = Presentation(pptx_file)
    results = []

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        # layout_name はステップ判定には使わない（判定はノート欄テキストのLLM意味分割）。
        # レイアウト名での判別は教員横断で約4割しか効かず廃止した。将来 LLM 結果の
        # 事後サニティチェック（レイアウト境界との突き合わせ）に使う余地として収集だけ残す。
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name or ""
        except AttributeError:
            # 取れないテンプレでも壊さない（AttributeError 限定＝実装バグは飲み込まない）
            layout_name = ""

        results.append({
            "slide_num": i,
            "title": title,
            "notes": notes,
            "layout_name": layout_name,
            "char_count": _count_reading_chars(notes),
            "pause_seconds": _count_pauses(notes),
        })

    return results


def analyze_narration(pptx_file, chars_per_min: int = 300) -> dict:
    """PPTXのナレーション文字数を集計し、推定尺を算出する。

    Args:
        pptx_file: PPTXファイルパスまたはファイルライクオブジェクト
        chars_per_min: 1分あたりの読み上げ文字数（デフォルト300）

    Returns:
        dict: 集計結果
            total_slides: 総スライド数
            slides_with_notes: ナレーション付きスライド数
            total_chars: 合計文字数（空白・改行除外）
            total_pause_seconds: 句読点ポーズ合計（秒）
            estimated_seconds: 推定尺（秒）= 読み上げ時間 + ポーズ
            slides: スライド別データのリスト
    """
    slides = extract_slide_notes(pptx_file)

    total_chars = sum(s["char_count"] for s in slides)
    total_pause = sum(s["pause_seconds"] for s in slides)
    slides_with_notes = sum(1 for s in slides if s["char_count"] > 0)
    reading_seconds = total_chars / chars_per_min * 60 if chars_per_min > 0 else 0
    estimated_seconds = reading_seconds + total_pause

    # 各スライドにも推定時間を追加
    for s in slides:
        reading = s["char_count"] / chars_per_min * 60 if chars_per_min > 0 else 0
        s["estimated_seconds"] = reading + s["pause_seconds"]

    return {
        "total_slides": len(slides),
        "slides_with_notes": slides_with_notes,
        "total_chars": total_chars,
        "total_pause_seconds": total_pause,
        "estimated_seconds": estimated_seconds,
        "slides": slides,
    }


def format_duration(seconds: float) -> str:
    """秒数を「○分○秒」形式にフォーマットする。"""
    minutes = int(seconds) // 60
    secs = int(seconds) % 60
    if minutes > 0:
        return f"{minutes}分{secs:02d}秒"
    return f"{secs}秒"


# ─────────────────────────────────────────────
# ステップ単位カウント
# ─────────────────────────────────────────────
# NeKo 大学ルール: 1ステップ 10〜17分（feedback_neko_video_length_constraint）。
# 1授業=最大4ステップ（大学運用ルール）。
# 境界の検出は step_segmenter.py が「ノート欄テキストの意味分割（Opus）」で行う。
# スライドのデザイン/レイアウトでの判別は教員横断で約4割しか効かない（実測）ため廃止。
# このモジュールは境界（とラベル）を受け取って集計するだけ＝関心分離。
STEP_TARGET_MIN_SEC = 10 * 60   # 600秒
STEP_TARGET_MAX_SEC = 17 * 60   # 1020秒


def group_slides_into_steps(slides, boundaries, labels=None,
                            target_min_sec=STEP_TARGET_MIN_SEC,
                            target_max_sec=STEP_TARGET_MAX_SEC):
    """slides を境界に従ってステップ単位に束ね、目標尺判定を付けて返す。

    Args:
        slides: analyze_narration の "slides"
        boundaries: 各ステップ開始スライド番号（昇順）。step_segmenter が返す or UIで上書き
        labels: 各ステップのラベル（boundaries と同順）。None/空要素は連番フォールバック

    Returns:
        list[dict]: step_num, step_label, first_title, slide_range, slide_count,
                    char_count, pause_seconds, estimated_seconds,
                    status("ok"|"over"|"under")
    """
    if not slides:
        return []
    slide_map = {s["slide_num"]: s for s in slides}
    all_nums = sorted(slide_map)
    bounds = sorted(n for n in set(boundaries) if n in slide_map)
    if not bounds:
        bounds = [all_nums[0]]

    steps = []
    for i, start in enumerate(bounds):
        end = bounds[i + 1] - 1 if i + 1 < len(bounds) else all_nums[-1]
        group = [slide_map[n] for n in all_nums if start <= n <= end]
        if not group:
            continue
        estimated = sum(s.get("estimated_seconds", 0.0) for s in group)
        if estimated > target_max_sec:
            status = "over"
        elif estimated < target_min_sec:
            status = "under"
        else:
            status = "ok"
        label = ""
        if labels and i < len(labels):
            label = (labels[i] or "").strip()
        steps.append({
            "step_num": i + 1,
            "step_label": label or f"STEP{i + 1}",
            "first_title": group[0].get("title", ""),
            "slide_range": (start, group[-1]["slide_num"]),
            "slide_count": len(group),
            "char_count": sum(s["char_count"] for s in group),
            "pause_seconds": sum(s["pause_seconds"] for s in group),
            "estimated_seconds": estimated,
            "status": status,
        })
    return steps
