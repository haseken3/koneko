"""PPTX 解析モジュール — 全スライドのタイトル・本文・ノート欄を抽出する。"""

from pptx import Presentation


def read_slides(pptx_file) -> list[dict]:
    """PPTX から全スライドのテキスト情報を抽出する。

    Args:
        pptx_file: PPTX ファイルパスまたはファイルライクオブジェクト

    Returns:
        list[dict]: 各スライドの情報
            slide_num (int): スライド番号（1始まり）
            title (str): スライドタイトル
            body (str): スライド本文（全テキストフレームを結合）
            notes (str): ノート欄テキスト
    """
    prs = Presentation(pptx_file)
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts = []

        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            text = sh.text_frame.text.strip()
            if not text:
                continue
            if not title:
                title = text.split("\n")[0][:120]
            body_parts.append(text)

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide_num": i,
            "title": title,
            "body": "\n".join(body_parts),
            "notes": notes,
        })

    return slides


def slides_summary(slides: list[dict]) -> dict:
    """全スライドの軽量サマリ（評価項目で全体参照する用）。

    Returns:
        dict: total_slides, titles (list of "S{num}: {title}")
    """
    return {
        "total_slides": len(slides),
        "titles": [f"S{s['slide_num']}: {s['title']}" for s in slides],
    }
